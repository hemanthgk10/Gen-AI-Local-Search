import PyPDF2
import docx
import pptx
import yaml
import os

class DocumentPreprocessor:
    @staticmethod
    def extract_text_from_pdf(file_path):
        """
        Extracts text from a PDF file. Joins all keys and values as a single string.
        """
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            return " ".join([page.extract_text() for page in reader.pages])

    @staticmethod
    def extract_text_from_docx(file_path):
        """
        Extracts text from a DOCX file. Joins all keys and values as a single string.
        """
        doc = docx.Document(file_path)
        return " ".join([para.text for para in doc.paragraphs])

    @staticmethod
    def extract_text_from_pptx(file_path):
        """
        Extracts text from a PPT file. Joins all keys and values as a single string.
        """
        ppt = pptx.Presentation(file_path)
        return " ".join([shape.text for slide in ppt.slides for shape in slide.shapes if shape.has_text_frame])
    
    @staticmethod
    def extract_text_from_yaml(file_path):
        """
        Extracts text from a YAML file. Joins all keys and values as a single string.
        """
        with open(file_path, "r") as file:
            data = yaml.safe_load(file)
            return DocumentPreprocessor.flatten_yaml(data)
        
    @staticmethod
    def flatten_yaml(data, indent=0):
        """
        Recursively flattens nested YAML content into a readable string.
        """
        lines = []
        if isinstance(data, dict):
            for key, value in data.items():
                lines.append(" " * indent + f"{key}:")
                lines.append(DocumentPreprocessor.flatten_yaml(value, indent=indent + 2))
        elif isinstance(data, list):
            for item in data:
                lines.append(DocumentPreprocessor.flatten_yaml(item, indent=indent + 2))
        else:
            lines.append(" " * indent + str(data))
        return "\n".join(lines)

    @staticmethod
    def preprocess_directory(directory_path):
        text_data = []
        for root, _, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                # print(f"File being processed: {file_path}")
                if file.endswith(".pdf"):
                    text_data.append(DocumentPreprocessor.extract_text_from_pdf(file_path))
                elif file.endswith(".docx"):
                    text_data.append(DocumentPreprocessor.extract_text_from_docx(file_path))
                elif file.endswith(".pptx"):
                    text_data.append(DocumentPreprocessor.extract_text_from_pptx(file_path))
                elif file.endswith((".yaml", ".yml")):
                    text_data.append(DocumentPreprocessor.extract_text_from_yaml(file_path))
        return text_data
